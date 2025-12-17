"""
Document Parser Module
Handles extraction of text and metadata from various document formats.
"""

import os
from datetime import datetime
from typing import Dict, Optional, Any
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl


class DocumentParser:
    """Parse various document formats and extract text and metadata."""
    
    def __init__(self):
        self.supported_formats = ['.pdf', '.docx', '.xlsx', '.pptx']
    
    def parse_document(self, file_path: str) -> Dict[str, Any]:
        """
        Parse document and return text content and metadata.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Dictionary containing text, metadata, and file info
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        # Get file metadata
        file_stats = os.stat(file_path)
        file_name = os.path.basename(file_path)
        
        result = {
            'file_name': file_name,
            'file_size': file_stats.st_size,
            'file_modified': datetime.fromtimestamp(file_stats.st_mtime).strftime('%Y-%m-%d %H:%M:%S'),
            'text_content': '',
            'metadata': {}
        }
        
        # Parse based on file type
        if file_extension == '.pdf':
            result.update(self._parse_pdf(file_path))
        elif file_extension == '.docx':
            result.update(self._parse_docx(file_path))
        elif file_extension == '.xlsx':
            result.update(self._parse_xlsx(file_path))
        elif file_extension == '.pptx':
            result.update(self._parse_pptx(file_path))
        
        return result
    
    def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """Extract text and metadata from PDF files."""
        text_content = []
        metadata = {}
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Extract metadata
                if pdf_reader.metadata:
                    metadata = {
                        'title': pdf_reader.metadata.get('/Title', ''),
                        'author': pdf_reader.metadata.get('/Author', ''),
                        'subject': pdf_reader.metadata.get('/Subject', ''),
                        'creator': pdf_reader.metadata.get('/Creator', ''),
                        'producer': pdf_reader.metadata.get('/Producer', ''),
                        'creation_date': pdf_reader.metadata.get('/CreationDate', ''),
                        'modification_date': pdf_reader.metadata.get('/ModDate', '')
                    }
                
                # Extract text from all pages
                for page in pdf_reader.pages:
                    text_content.append(page.extract_text())
        
        except Exception as e:
            raise Exception(f"Error parsing PDF: {str(e)}")
        
        return {
            'text_content': '\n'.join(text_content),
            'metadata': metadata
        }
    
    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        """Extract text and metadata from Word documents."""
        text_content = []
        metadata = {}
        
        try:
            doc = Document(file_path)
            
            # Extract core properties
            core_properties = doc.core_properties
            metadata = {
                'title': core_properties.title or '',
                'author': core_properties.author or '',
                'subject': core_properties.subject or '',
                'keywords': core_properties.keywords or '',
                'comments': core_properties.comments or '',
                'created': core_properties.created.strftime('%Y-%m-%d %H:%M:%S') if core_properties.created else '',
                'modified': core_properties.modified.strftime('%Y-%m-%d %H:%M:%S') if core_properties.modified else '',
                'last_modified_by': core_properties.last_modified_by or ''
            }
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join([cell.text for cell in row.cells])
                    if row_text.strip():
                        text_content.append(row_text)
        
        except Exception as e:
            raise Exception(f"Error parsing DOCX: {str(e)}")
        
        return {
            'text_content': '\n'.join(text_content),
            'metadata': metadata
        }
    
    def _parse_xlsx(self, file_path: str) -> Dict[str, Any]:
        """Extract text and metadata from Excel spreadsheets."""
        text_content = []
        metadata = {}
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            
            # Extract properties
            props = workbook.properties
            metadata = {
                'title': props.title or '',
                'creator': props.creator or '',
                'subject': props.subject or '',
                'keywords': props.keywords or '',
                'description': props.description or '',
                'created': props.created.strftime('%Y-%m-%d %H:%M:%S') if props.created else '',
                'modified': props.modified.strftime('%Y-%m-%d %H:%M:%S') if props.modified else '',
                'last_modified_by': props.lastModifiedBy or ''
            }
            
            # Extract text from all sheets
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content.append(f"Sheet: {sheet_name}")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join([str(cell) if cell is not None else '' for cell in row])
                    if row_text.strip():
                        text_content.append(row_text)
        
        except Exception as e:
            raise Exception(f"Error parsing XLSX: {str(e)}")
        
        return {
            'text_content': '\n'.join(text_content),
            'metadata': metadata
        }
    
    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """Extract text and metadata from PowerPoint presentations."""
        text_content = []
        metadata = {}
        
        try:
            presentation = Presentation(file_path)
            
            # Extract core properties
            core_properties = presentation.core_properties
            metadata = {
                'title': core_properties.title or '',
                'author': core_properties.author or '',
                'subject': core_properties.subject or '',
                'keywords': core_properties.keywords or '',
                'comments': core_properties.comments or '',
                'created': core_properties.created.strftime('%Y-%m-%d %H:%M:%S') if core_properties.created else '',
                'modified': core_properties.modified.strftime('%Y-%m-%d %H:%M:%S') if core_properties.modified else '',
                'last_modified_by': core_properties.last_modified_by or ''
            }
            
            # Extract text from all slides
            for i, slide in enumerate(presentation.slides, 1):
                text_content.append(f"Slide {i}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text)
        
        except Exception as e:
            raise Exception(f"Error parsing PPTX: {str(e)}")
        
        return {
            'text_content': '\n'.join(text_content),
            'metadata': metadata
        }
