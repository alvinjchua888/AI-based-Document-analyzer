"""
Tests for document_parser module
"""

import pytest
import os
from document_parser import DocumentParser
from docx import Document
from pptx import Presentation
import openpyxl


@pytest.fixture
def parser():
    """Create a DocumentParser instance."""
    return DocumentParser()


@pytest.fixture
def sample_docx(tmp_path):
    """Create a sample Word document."""
    doc_path = tmp_path / "test.docx"
    doc = Document()
    doc.core_properties.title = "Test Document"
    doc.core_properties.author = "Test Author"
    doc.add_paragraph("This is a test document.")
    doc.save(str(doc_path))
    return str(doc_path)


@pytest.fixture
def sample_pptx(tmp_path):
    """Create a sample PowerPoint presentation."""
    pptx_path = tmp_path / "test.pptx"
    prs = Presentation()
    prs.core_properties.title = "Test Presentation"
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Slide"
    prs.save(str(pptx_path))
    return str(pptx_path)


@pytest.fixture
def sample_xlsx(tmp_path):
    """Create a sample Excel spreadsheet."""
    xlsx_path = tmp_path / "test.xlsx"
    wb = openpyxl.Workbook()
    wb.properties.title = "Test Spreadsheet"
    ws = wb.active
    ws['A1'] = "Test"
    ws['B1'] = "Data"
    wb.save(str(xlsx_path))
    return str(xlsx_path)


def test_parser_initialization(parser):
    """Test that parser initializes with correct supported formats."""
    assert parser.supported_formats == ['.pdf', '.docx', '.xlsx', '.pptx']


def test_parse_docx(parser, sample_docx):
    """Test parsing Word document."""
    result = parser.parse_document(sample_docx)
    
    assert 'file_name' in result
    assert 'text_content' in result
    assert 'metadata' in result
    assert result['file_name'] == 'test.docx'
    assert 'This is a test document' in result['text_content']
    assert result['metadata']['title'] == 'Test Document'
    assert result['metadata']['author'] == 'Test Author'


def test_parse_pptx(parser, sample_pptx):
    """Test parsing PowerPoint presentation."""
    result = parser.parse_document(sample_pptx)
    
    assert 'file_name' in result
    assert 'text_content' in result
    assert 'metadata' in result
    assert result['file_name'] == 'test.pptx'
    assert 'Test Slide' in result['text_content']
    assert result['metadata']['title'] == 'Test Presentation'


def test_parse_xlsx(parser, sample_xlsx):
    """Test parsing Excel spreadsheet."""
    result = parser.parse_document(sample_xlsx)
    
    assert 'file_name' in result
    assert 'text_content' in result
    assert 'metadata' in result
    assert result['file_name'] == 'test.xlsx'
    assert 'Test' in result['text_content']
    assert result['metadata']['title'] == 'Test Spreadsheet'


def test_unsupported_format(parser, tmp_path):
    """Test that unsupported formats raise ValueError."""
    txt_file = tmp_path / "test.txt"
    txt_file.write_text("Test content")
    
    with pytest.raises(ValueError) as excinfo:
        parser.parse_document(str(txt_file))
    
    assert "Unsupported file format" in str(excinfo.value)


def test_parse_nonexistent_file(parser):
    """Test parsing non-existent file."""
    with pytest.raises(Exception):
        parser.parse_document("/nonexistent/file.docx")
